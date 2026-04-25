"""
pptx_panjit.py — PANJIT 品牌母片可編輯 PPTX 生成工具庫

用途
----
Code-first 建立符合 PANJIT 品牌母片規範的 PowerPoint 報告。固定的是封面、
header、logo、頁尾、頁碼與內容安全區；中間內容頁仍由 python-pptx 原生
自由排版，產出 100% 可編輯的 .pptx，與 Marp（HTML 渲染）互補。

設計原則
--------
1. 尺寸全部用 Inches()/Pt()/RGBColor()，不直接操作 EMU 整數
2. 顏色輸入統一用 6 位 hex string (e.g. "1F4E79")，內部轉 RGBColor
3. 座標來自既有 PANJIT 簡報版面分析（見 constants 區塊）
4. 每一頁自動加品牌固定元素（版權 + 頁碼）
5. 品牌母片 presentation 是殼，add_xxx 函式作用在 slide 物件上

API 速查
--------
class PanjitBrandMasterPresentation
    __init__(template_path=None)
    add_title_slide(title, subtitle, department, owner, date, version)
    add_section_divider(section_title)
    add_master_content_slide(title, section_label="") -> slide
    get_content_safe_area() -> dict
    save(output_path)

functions (on slide)
    add_table(slide, data, col_widths, top, left, width, header_color, alt_row, font_size, caption)
    add_image(slide, image_path, left, top, width, height, caption)
    add_bullets(slide, items, top, left, width, font_size, title)
    add_two_column(slide, left_content_fn, right_content_fn, ratio)
    add_text_box(slide, text, left, top, width, height, font_size, bold, color, bg_color, border_left_color)

utilities
    set_cell_style(cell, bg_color, font_color, font_size, bold, align)
    apply_table_style(table, header_color, alt_row_color, ...)

CLI
---
    python pptx_panjit.py demo output.pptx
    python pptx_panjit.py dump-spec brand_spec.json
    python pptx_panjit.py demo-with-spec brand_spec.json output.pptx
    python pptx_panjit.py info template.pptx
"""

from __future__ import annotations

import argparse
import json
import posixpath
import re
import subprocess
import sys
import xml.etree.ElementTree as ET
import zipfile
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Callable, Optional, Sequence

# Avoid local `shared/tools/pptx/` shadowing the third-party `python-pptx` package.
_THIS_DIR = Path(__file__).resolve().parent
sys.path = [p for p in sys.path if Path(p).resolve() != _THIS_DIR]

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ============================================================================
# 版面常數（inches）— 來自對 MBU1 Non-AU 新增銀膠勤凱供應商.pptx 的座標分析
# ============================================================================

SLIDE_WIDTH = 13.33
SLIDE_HEIGHT = 7.50

# 標題列
TITLE_LEFT = 0.59
TITLE_TOP = 0.12
TITLE_WIDTH = 10.12
TITLE_HEIGHT = 0.80

# 次標題/節標籤（放在標題下方，字級較小）
SECTION_LABEL_LEFT = 0.59
SECTION_LABEL_TOP = 0.94
SECTION_LABEL_WIDTH = 10.12
SECTION_LABEL_HEIGHT = 0.30

# 內容區
CONTENT_TOP = 1.35
CONTENT_LEFT = 0.42
CONTENT_WIDTH = 12.5
CONTENT_BOTTOM = 7.00  # 避開頁尾
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP  # ~5.65

# 頁尾
FOOTER_LEFT = 0.22
FOOTER_TOP = 7.19
FOOTER_WIDTH = 4.82
FOOTER_HEIGHT = 0.28
FOOTER_TEXT = "Copyright© 2020 PANJIT International Inc. All rights reserved."

# 頁碼位置（右下）
PAGE_NUM_LEFT = 12.50
PAGE_NUM_TOP = 7.19
PAGE_NUM_WIDTH = 0.60
PAGE_NUM_HEIGHT = 0.28

# 內容頁固定品牌元素（來自 PANJIT 範例母片的實測座標）
HEADER_LINE_LEFT = 0.65
HEADER_LINE_TOP = 0.80
HEADER_LINE_WIDTH = 10.65
HEADER_LINE_HEIGHT = 0.03

CONTENT_LOGO_LEFT = 11.34
CONTENT_LOGO_TOP = 0.23
CONTENT_LOGO_WIDTH = 1.51
CONTENT_LOGO_HEIGHT = 0.61

# ============================================================================
# 顏色常數（semicon.css 色盤）
# ============================================================================

COLOR_PRIMARY = "1F4E79"       # 深藍，標題/表頭
COLOR_SECONDARY = "D6E4F0"     # 淺藍，次要背景
COLOR_ACCENT = "F4B084"        # 橘，強調色
COLOR_SUCCESS = "C6EFCE"       # 綠，通過/OK
COLOR_WARNING = "FFEB9C"       # 黃，警告
COLOR_DANGER = "FFC7CE"        # 紅，失敗/NG
COLOR_ALT_ROW = "F2F2F2"       # 淺灰，交替行
COLOR_WHITE = "FFFFFF"
COLOR_TEXT = "1F2937"          # 深灰，主文字
COLOR_MUTED = "6B7280"         # 中灰，次要文字

# ============================================================================
# 字型常數
# ============================================================================

FONT_CJK = "Microsoft JhengHei"
FONT_LATIN = "Calibri"

# 字級
FONT_SIZE_COVER_TITLE = 36
FONT_SIZE_COVER_SUB = 20
FONT_SIZE_COVER_META = 14
FONT_SIZE_SECTION = 40
FONT_SIZE_TITLE = 20
FONT_SIZE_SUBTITLE = 14
FONT_SIZE_BODY = 12
FONT_SIZE_TABLE_HEADER = 10.5
FONT_SIZE_TABLE_DATA = 10
FONT_SIZE_CAPTION = 10
FONT_SIZE_FOOTER = 10


@dataclass
class RectSpec:
    left: float
    top: float
    width: float
    height: float


@dataclass
class BrandSpec:
    slide_width: float = SLIDE_WIDTH
    slide_height: float = SLIDE_HEIGHT
    title_area: RectSpec = field(
        default_factory=lambda: RectSpec(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    )
    section_label_area: RectSpec = field(
        default_factory=lambda: RectSpec(
            SECTION_LABEL_LEFT, SECTION_LABEL_TOP, SECTION_LABEL_WIDTH, SECTION_LABEL_HEIGHT
        )
    )
    content_area: RectSpec = field(
        default_factory=lambda: RectSpec(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    )
    footer_area: RectSpec = field(
        default_factory=lambda: RectSpec(FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH, FOOTER_HEIGHT)
    )
    page_num_area: RectSpec = field(
        default_factory=lambda: RectSpec(PAGE_NUM_LEFT, PAGE_NUM_TOP, PAGE_NUM_WIDTH, PAGE_NUM_HEIGHT)
    )
    header_line_area: RectSpec = field(
        default_factory=lambda: RectSpec(
            HEADER_LINE_LEFT, HEADER_LINE_TOP, HEADER_LINE_WIDTH, HEADER_LINE_HEIGHT
        )
    )
    content_logo_area: RectSpec = field(
        default_factory=lambda: RectSpec(
            CONTENT_LOGO_LEFT, CONTENT_LOGO_TOP, CONTENT_LOGO_WIDTH, CONTENT_LOGO_HEIGHT
        )
    )
    footer_text: str = FOOTER_TEXT
    color_primary: str = COLOR_PRIMARY
    color_secondary: str = COLOR_SECONDARY
    color_accent: str = COLOR_ACCENT
    color_text: str = COLOR_TEXT
    color_muted: str = COLOR_MUTED
    color_white: str = COLOR_WHITE
    font_cjk: str = FONT_CJK
    font_latin: str = FONT_LATIN
    content_logo_path: Optional[str] = None
    cover_logo_path: Optional[str] = None

    @staticmethod
    def _rect_from_dict(raw: dict, fallback: RectSpec) -> RectSpec:
        return RectSpec(
            left=float(raw.get("left", fallback.left)),
            top=float(raw.get("top", fallback.top)),
            width=float(raw.get("width", fallback.width)),
            height=float(raw.get("height", fallback.height)),
        )

    @classmethod
    def from_dict(cls, data: dict) -> "BrandSpec":
        spec = cls()
        spec.slide_width = float(data.get("slide_width", spec.slide_width))
        spec.slide_height = float(data.get("slide_height", spec.slide_height))
        spec.title_area = cls._rect_from_dict(data.get("title_area", {}), spec.title_area)
        spec.section_label_area = cls._rect_from_dict(
            data.get("section_label_area", {}), spec.section_label_area
        )
        spec.content_area = cls._rect_from_dict(data.get("content_area", {}), spec.content_area)
        spec.footer_area = cls._rect_from_dict(data.get("footer_area", {}), spec.footer_area)
        spec.page_num_area = cls._rect_from_dict(data.get("page_num_area", {}), spec.page_num_area)
        spec.header_line_area = cls._rect_from_dict(
            data.get("header_line_area", {}), spec.header_line_area
        )
        spec.content_logo_area = cls._rect_from_dict(
            data.get("content_logo_area", {}), spec.content_logo_area
        )
        spec.footer_text = data.get("footer_text", spec.footer_text)
        spec.color_primary = data.get("color_primary", spec.color_primary)
        spec.color_secondary = data.get("color_secondary", spec.color_secondary)
        spec.color_accent = data.get("color_accent", spec.color_accent)
        spec.color_text = data.get("color_text", spec.color_text)
        spec.color_muted = data.get("color_muted", spec.color_muted)
        spec.color_white = data.get("color_white", spec.color_white)
        spec.font_cjk = data.get("font_cjk", spec.font_cjk)
        spec.font_latin = data.get("font_latin", spec.font_latin)
        spec.content_logo_path = data.get("content_logo_path", spec.content_logo_path)
        spec.cover_logo_path = data.get("cover_logo_path", spec.cover_logo_path)
        return spec

    def to_dict(self) -> dict:
        return {
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
            "title_area": asdict(self.title_area),
            "section_label_area": asdict(self.section_label_area),
            "content_area": asdict(self.content_area),
            "footer_area": asdict(self.footer_area),
            "page_num_area": asdict(self.page_num_area),
            "header_line_area": asdict(self.header_line_area),
            "content_logo_area": asdict(self.content_logo_area),
            "footer_text": self.footer_text,
            "color_primary": self.color_primary,
            "color_secondary": self.color_secondary,
            "color_accent": self.color_accent,
            "color_text": self.color_text,
            "color_muted": self.color_muted,
            "color_white": self.color_white,
            "font_cjk": self.font_cjk,
            "font_latin": self.font_latin,
            "content_logo_path": self.content_logo_path,
            "cover_logo_path": self.cover_logo_path,
        }


DEFAULT_BRAND_SPEC = BrandSpec()
_PML_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
_REPO_ROOT = _THIS_DIR.parent.parent
_PANJIT_BRAND_SPEC_PATH = _REPO_ROOT / ".claude" / "skills-on-demand" / "pptx-brand-master" / "brand_spec.panjit.json"


# ============================================================================
# 顏色轉換小工具
# ============================================================================

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """'1F4E79' -> RGBColor. 接受 '#1F4E79' 或 '1F4E79'."""
    if not hex_color:
        return None
    h = hex_color.lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Invalid hex color: {hex_color}")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_run_font(run, *, size=None, bold=None, color_hex=None,
                  name_cjk=FONT_CJK, name_latin=FONT_LATIN):
    """套用字型設定到 run。size 為 pt 數值。"""
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color_hex is not None:
        run.font.color.rgb = _hex_to_rgb(color_hex)
    # 同時指定 latin 與 east asian 字型
    run.font.name = name_latin
    rPr = run._r.get_or_add_rPr()
    # 東亞字型（中文）— 透過 rFonts 設定
    from pptx.oxml.ns import qn
    rFonts = rPr.find(qn("a:rFonts"))
    if rFonts is None:
        rFonts = rPr.makeelement(qn("a:rFonts"), {})
        rPr.insert(0, rFonts)
    rFonts.set("eastAsia", name_cjk)
    rFonts.set("ascii", name_latin)
    rFonts.set("hAnsi", name_latin)


def _set_shape_fill(shape, hex_color: str):
    """實心填色"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _set_shape_no_line(shape):
    """移除外框線"""
    shape.line.fill.background()


def _set_slide_background(slide, hex_color: str) -> None:
    """設定 slide 背景色（在 slide 層，不產生 shape）。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _normalize_cli_path(path: str) -> str:
    """
    Normalize paths passed through PowerShell/Git Bash.

    Examples:
    - /d/ai-office/foo.pptx -> D:\\ai-office\\foo.pptx
    - D:\\ai-office\\foo.pptx -> unchanged
    """
    if not path:
        return path
    m = re.match(r"^/([a-zA-Z])/(.+)$", path)
    if m:
        drive = m.group(1).upper()
        tail = m.group(2).replace("/", "\\")
        return f"{drive}:\\{tail}"
    return path


def load_brand_spec(path: str) -> BrandSpec:
    raw = json.loads(Path(_normalize_cli_path(path)).read_text(encoding="utf-8"))
    return BrandSpec.from_dict(raw)


def save_brand_spec(path: str, brand_spec: BrandSpec) -> None:
    out = Path(_normalize_cli_path(path))
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(
        json.dumps(brand_spec.to_dict(), ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def get_active_brand_spec() -> BrandSpec:
    if _PANJIT_BRAND_SPEC_PATH.exists():
        return load_brand_spec(str(_PANJIT_BRAND_SPEC_PATH))
    return BrandSpec()


def _bind_slide_brand_context(slide, brand_spec: BrandSpec) -> None:
    setattr(slide, "_brand_spec", brand_spec)


def _get_slide_brand_spec(slide) -> BrandSpec:
    return getattr(slide, "_brand_spec", get_active_brand_spec())


# ============================================================================
# PanjitBrandMasterPresentation — 主類別
# ============================================================================

class PanjitBrandMasterPresentation:
    """
    PANJIT 品牌母片可編輯簡報。

    用法::

        prs = PanjitBrandMasterPresentation()
        prs.add_title_slide("月度品質報告", "2026 Q1",
                            department="QA", owner="林志明", date="2026-04-21")
        prs.add_section_divider("可靠性測試結果")
        slide = prs.add_master_content_slide("AEC-Q101 結果", section_label="Phase 1")
        add_table(slide, [["項目","規格","結果"], ["HTRL","1000h","PASS"]])
        prs.save("report.pptx")
    """

    def __init__(
        self,
        template_path: Optional[str] = None,
        brand_spec: Optional[BrandSpec] = None,
        brand_spec_path: Optional[str] = None,
    ):
        if brand_spec_path:
            self.brand_spec = load_brand_spec(brand_spec_path)
        else:
            self.brand_spec = brand_spec or get_active_brand_spec()

        normalized_template_path = _normalize_cli_path(template_path) if template_path else None
        if normalized_template_path and Path(normalized_template_path).exists():
            self.prs = Presentation(normalized_template_path)
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(self.brand_spec.slide_width)
            self.prs.slide_height = Inches(self.brand_spec.slide_height)

        # 空白佈局索引（layout 6 通常是空白；若模板不同則退回 layout 0）
        blank_layouts = [l for l in self.prs.slide_layouts if l.name.lower() == "blank"]
        self._blank_layout = blank_layouts[0] if blank_layouts else self.prs.slide_layouts[
            min(6, len(self.prs.slide_layouts) - 1)
        ]

    @classmethod
    def from_brand_spec_file(
        cls,
        spec_path: str,
        template_path: Optional[str] = None,
    ) -> "PanjitBrandMasterPresentation":
        return cls(template_path=template_path, brand_spec_path=spec_path)

    # ------------------------------------------------------------------
    # 封面頁
    # ------------------------------------------------------------------
    def add_title_slide(self, title: str, subtitle: str = "",
                        department: str = "", owner: str = "",
                        date: str = "", version: str = "V0") -> object:
        slide = self.prs.slides.add_slide(self._blank_layout)
        spec = self.brand_spec
        _bind_slide_brand_context(slide, spec)

        # 深藍全頁背景（用 slide.background，不產生 shape）
        _set_slide_background(slide, spec.color_primary)

        # 標題橫條（放在頁面最底，與其他元素不重疊）
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.0), Inches(6.90),
            Inches(spec.slide_width), Inches(0.08)
        )
        _set_shape_fill(accent_bar, spec.color_accent)
        _set_shape_no_line(accent_bar)

        self._add_cover_logo(slide)

        # 大標題
        title_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.3), Inches(spec.slide_width - 2.0), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title or ""
        _set_run_font(run, size=FONT_SIZE_COVER_TITLE, bold=True, color_hex=spec.color_white)

        # 副標題
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(3.6), Inches(spec.slide_width - 2.0), Inches(0.8)
            )
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            tf2.margin_left = tf2.margin_right = 0
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = subtitle
            _set_run_font(r2, size=FONT_SIZE_COVER_SUB, bold=False, color_hex=spec.color_secondary)

        # Metadata 區（部門 / 負責人 / 日期 / 版本）
        meta_lines = []
        if department:
            meta_lines.append(f"Department: {department}")
        if owner:
            meta_lines.append(f"Owner: {owner}")
        if date:
            meta_lines.append(f"Date: {date}")
        if version:
            meta_lines.append(f"Version: {version}")

        if meta_lines:
            meta_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(5.2), Inches(spec.slide_width - 2.0), Inches(1.5)
            )
            tf3 = meta_box.text_frame
            tf3.word_wrap = True
            tf3.margin_left = tf3.margin_right = 0
            for idx, line in enumerate(meta_lines):
                p3 = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
                p3.alignment = PP_ALIGN.LEFT
                r3 = p3.add_run()
                r3.text = line
                _set_run_font(r3, size=FONT_SIZE_COVER_META, bold=False, color_hex=spec.color_white)
                p3.space_after = Pt(2)

        return slide

    # ------------------------------------------------------------------
    # 章節分隔頁
    # ------------------------------------------------------------------
    def add_section_divider(self, section_title: str) -> object:
        slide = self.prs.slides.add_slide(self._blank_layout)
        spec = self.brand_spec
        _bind_slide_brand_context(slide, spec)

        # 淺藍全頁背景（用 slide.background，不產生 shape）
        _set_slide_background(slide, spec.color_secondary)

        # 左側深藍色條（橫向 0.50"–0.65"，與文字 0.95" 起不重疊）
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(2.8),
            Inches(0.15), Inches(2.0)
        )
        _set_shape_fill(bar, spec.color_primary)
        _set_shape_no_line(bar)

        # 大字章節標題
        title_box = slide.shapes.add_textbox(
            Inches(0.95), Inches(3.1),
            Inches(SLIDE_WIDTH - 2.0), Inches(1.4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = section_title or ""
        _set_run_font(r, size=FONT_SIZE_SECTION, bold=True, color_hex=spec.color_primary)

        self._add_content_logo(slide)

        self._add_footer(slide)
        return slide

    # ------------------------------------------------------------------
    # 內容頁（回傳 slide 讓呼叫者加內容）
    # ------------------------------------------------------------------
    def add_master_content_slide(self, title: str, section_label: str = "") -> object:
        slide = self.prs.slides.add_slide(self._blank_layout)
        spec = self.brand_spec
        _bind_slide_brand_context(slide, spec)

        # 標題列
        title_box = slide.shapes.add_textbox(
            Inches(spec.title_area.left), Inches(spec.title_area.top),
            Inches(spec.title_area.width), Inches(spec.title_area.height)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title or ""
        _set_run_font(r, size=FONT_SIZE_TITLE, bold=True, color_hex=spec.color_primary)

        # 節標籤（選填，放在標題下方）
        label_bottom = spec.title_area.top + spec.title_area.height
        if section_label:
            label_top = spec.title_area.top + spec.title_area.height + 0.02
            label_h = spec.section_label_area.height
            label_box = slide.shapes.add_textbox(
                Inches(spec.section_label_area.left), Inches(label_top),
                Inches(spec.section_label_area.width), Inches(label_h)
            )
            tf2 = label_box.text_frame
            tf2.word_wrap = True
            tf2.margin_left = tf2.margin_right = 0
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = section_label
            _set_run_font(r2, size=FONT_SIZE_SUBTITLE, bold=False, color_hex=spec.color_muted)
            label_bottom = label_top + label_h

        # 分隔線（放在標題/節標籤之下）
        line_top = max(label_bottom + 0.02, spec.header_line_area.top)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(spec.header_line_area.left), Inches(line_top),
            Inches(spec.header_line_area.width), Inches(spec.header_line_area.height)
        )
        _set_shape_fill(line, spec.color_primary)
        _set_shape_no_line(line)

        self._add_content_logo(slide)

        self._add_footer(slide)
        return slide

    def add_content_slide(self, title: str, section_label: str = "") -> object:
        """
        Backward-compatible alias.

        舊名稱保留，但語義上這是「套品牌母片的內容頁骨架」，
        不是固定內容模板。
        """
        return self.add_master_content_slide(title, section_label=section_label)

    def get_content_safe_area(self) -> dict:
        """回傳品牌母片中的標題區、內容安全區與頁尾區（inches）。"""
        spec = self.brand_spec
        return {
            "title": {
                "left": spec.title_area.left,
                "top": spec.title_area.top,
                "width": spec.title_area.width,
                "height": spec.title_area.height,
            },
            "section_label": {
                "left": spec.section_label_area.left,
                "top": spec.section_label_area.top,
                "width": spec.section_label_area.width,
                "height": spec.section_label_area.height,
            },
            "content": {
                "left": spec.content_area.left,
                "top": spec.content_area.top,
                "width": spec.content_area.width,
                "height": spec.content_area.height,
                "bottom": spec.content_area.top + spec.content_area.height,
            },
            "header_line": {
                "left": spec.header_line_area.left,
                "top": spec.header_line_area.top,
                "width": spec.header_line_area.width,
                "height": spec.header_line_area.height,
            },
            "content_logo": {
                "left": spec.content_logo_area.left,
                "top": spec.content_logo_area.top,
                "width": spec.content_logo_area.width,
                "height": spec.content_logo_area.height,
            },
            "footer": {
                "left": spec.footer_area.left,
                "top": spec.footer_area.top,
                "width": spec.footer_area.width,
                "height": spec.footer_area.height,
            },
            "page_number": {
                "left": spec.page_num_area.left,
                "top": spec.page_num_area.top,
                "width": spec.page_num_area.width,
                "height": spec.page_num_area.height,
            },
        }

    # ------------------------------------------------------------------
    # 頁尾（私有）
    # ------------------------------------------------------------------
    def _add_logo(self, slide, path: Optional[str], area: RectSpec) -> None:
        if not path:
            return
        logo = Path(path)
        if not logo.exists():
            return
        slide.shapes.add_picture(
            str(logo),
            Inches(area.left),
            Inches(area.top),
            width=Inches(area.width),
            height=Inches(area.height),
        )

    def _add_content_logo(self, slide) -> None:
        self._add_logo(slide, self.brand_spec.content_logo_path, self.brand_spec.content_logo_area)

    def _add_cover_logo(self, slide) -> None:
        path = self.brand_spec.cover_logo_path or self.brand_spec.content_logo_path
        self._add_logo(slide, path, self.brand_spec.content_logo_area)

    def _add_footer(self, slide) -> None:
        spec = self.brand_spec
        footer_top = min(
            spec.footer_area.top,
            max(0.0, spec.slide_height - spec.footer_area.height),
        )
        page_num_top = min(
            spec.page_num_area.top,
            max(0.0, spec.slide_height - spec.page_num_area.height),
        )
        # 版權
        foot = slide.shapes.add_textbox(
            Inches(spec.footer_area.left), Inches(footer_top),
            Inches(spec.footer_area.width), Inches(spec.footer_area.height)
        )
        tf = foot.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = spec.footer_text
        _set_run_font(r, size=FONT_SIZE_FOOTER, bold=False, color_hex=spec.color_muted)

        # 頁碼
        num = slide.shapes.add_textbox(
            Inches(spec.page_num_area.left), Inches(page_num_top),
            Inches(spec.page_num_area.width), Inches(spec.page_num_area.height)
        )
        tf2 = num.text_frame
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        # 頁碼 = 當前 slide index + 1
        r2.text = str(len(self.prs.slides))
        _set_run_font(r2, size=FONT_SIZE_FOOTER, bold=False, color_hex=spec.color_muted)

    # ------------------------------------------------------------------
    def save(self, output_path: str) -> None:
        out = Path(_normalize_cli_path(output_path))
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))


# ============================================================================
# 格式化輔助
# ============================================================================

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def set_cell_style(cell, *, bg_color: Optional[str] = None,
                   font_color: Optional[str] = None,
                   font_size: Optional[float] = None,
                   bold: Optional[bool] = None,
                   align: str = "left") -> None:
    """設定單一 table cell 的樣式。"""
    if bg_color is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb(bg_color)

    tf = cell.text_frame
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for p in tf.paragraphs:
        p.alignment = _ALIGN_MAP.get(align, PP_ALIGN.LEFT)
        for r in p.runs:
            if font_size is not None:
                r.font.size = Pt(font_size)
            if font_color is not None:
                r.font.color.rgb = _hex_to_rgb(font_color)
            if bold is not None:
                r.font.bold = bold
            # 套字型
            _set_run_font(
                r,
                size=font_size,
                bold=bold,
                color_hex=font_color,
            )


def apply_table_style(table, *,
                      header_color: str = COLOR_PRIMARY,
                      alt_row_color: Optional[str] = COLOR_ALT_ROW,
                      header_font_color: str = COLOR_WHITE,
                      header_font_size: float = FONT_SIZE_TABLE_HEADER,
                      data_font_size: float = FONT_SIZE_TABLE_DATA,
                      data_font_color: str = COLOR_TEXT) -> None:
    """整表套樣式。跳過合併儲存格（非起始格）。"""
    n_rows = len(table.rows)
    n_cols = len(table.columns)

    for r_idx in range(n_rows):
        for c_idx in range(n_cols):
            cell = table.cell(r_idx, c_idx)
            # 跳過合併儲存格的延伸部分
            if cell.is_spanned:
                continue

            if r_idx == 0:
                # header
                set_cell_style(
                    cell,
                    bg_color=header_color,
                    font_color=header_font_color,
                    font_size=header_font_size,
                    bold=True,
                    align="center",
                )
            else:
                bg = None
                if alt_row_color and (r_idx % 2 == 0):
                    bg = alt_row_color
                set_cell_style(
                    cell,
                    bg_color=bg,
                    font_color=data_font_color,
                    font_size=data_font_size,
                    bold=False,
                    align="left",
                )


# ============================================================================
# 元素建構函式（作用在 slide 上）
# ============================================================================

def add_table(slide, data: Sequence[Sequence],
              col_widths: Optional[Sequence[float]] = None,
              top: Optional[float] = None,
              left: Optional[float] = None,
              width: Optional[float] = None,
              header_color: str = COLOR_PRIMARY,
              alt_row: bool = True,
              font_size: float = FONT_SIZE_TABLE_DATA,
              caption: Optional[str] = None):
    """
    加入表格到 slide。

    Args:
        data: 二維陣列，第一行是 header
        col_widths: 各欄寬度（inches），None 則均分
        top/left/width: 位置與總寬（inches），None 用內容區預設
        header_color: 表頭底色 hex
        alt_row: 偶數資料行交替淺灰
        font_size: 資料行字級
        caption: 表格下方說明文字
    """
    if not data or not data[0]:
        raise ValueError("data must be a non-empty 2D list")

    spec = _get_slide_brand_spec(slide)
    n_rows = len(data)
    n_cols = len(data[0])

    # 位置預設
    _left = left if left is not None else spec.content_area.left
    _top = top if top is not None else spec.content_area.top
    _width = width if width is not None else spec.content_area.width

    # 估算行高：資料行約 0.30"，header 0.35"
    row_h = 0.30
    header_h = 0.35
    total_h = header_h + row_h * (n_rows - 1)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(_left), Inches(_top),
        Inches(_width), Inches(total_h)
    )
    table = table_shape.table

    # 列高
    table.rows[0].height = Inches(header_h)
    for i in range(1, n_rows):
        table.rows[i].height = Inches(row_h)

    # 欄寬
    if col_widths is None:
        col_widths = [_width / n_cols] * n_cols
    else:
        if len(col_widths) != n_cols:
            raise ValueError(
                f"col_widths length {len(col_widths)} != n_cols {n_cols}"
            )
    for idx, w in enumerate(col_widths):
        table.columns[idx].width = Inches(w)

    # 填值
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            if c_idx >= n_cols:
                continue
            cell = table.cell(r_idx, c_idx)
            if cell.is_spanned:
                continue
            cell.text = "" if val is None else str(val)

    # 套樣式
    apply_table_style(
        table,
        header_color=header_color,
        alt_row_color=COLOR_ALT_ROW if alt_row else None,
        data_font_size=font_size,
    )

    # caption（表格下方）
    if caption:
        cap_top = _top + total_h + 0.05
        cap_box = slide.shapes.add_textbox(
            Inches(_left), Inches(cap_top),
            Inches(_width), Inches(0.28)
        )
        tf = cap_box.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = caption
        _set_run_font(r, size=FONT_SIZE_CAPTION, bold=False, color_hex=COLOR_MUTED)

    return table_shape


def add_image(slide, image_path: str,
              left: Optional[float] = None,
              top: Optional[float] = None,
              width: Optional[float] = None,
              height: Optional[float] = None,
              caption: Optional[str] = None):
    """
    嵌入圖片。位置未指定時置中於內容區。

    Args:
        image_path: 圖檔路徑（絕對建議）
        left/top/width/height: inches
        caption: 圖片下方說明
    """
    img_path = Path(_normalize_cli_path(image_path))
    if not img_path.exists():
        raise FileNotFoundError(f"Image not found: {image_path}")

    spec = _get_slide_brand_spec(slide)
    # 位置預設：置中於內容區
    if width is None and height is None:
        width = min(8.0, spec.content_area.width - 1.0)
    if left is None:
        if width is not None:
            left = spec.content_area.left + (spec.content_area.width - width) / 2
        else:
            left = spec.content_area.left
    if top is None:
        top = spec.content_area.top + 0.2

    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)

    pic = slide.shapes.add_picture(
        str(img_path), Inches(left), Inches(top), **kwargs
    )

    if caption:
        # 用 pic 實際尺寸算 caption 位置
        cap_top = Emu(pic.top + pic.height).inches + 0.05
        cap_left = Emu(pic.left).inches
        cap_w = Emu(pic.width).inches
        cap_box = slide.shapes.add_textbox(
            Inches(cap_left), Inches(cap_top),
            Inches(cap_w), Inches(0.28)
        )
        tf = cap_box.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        _set_run_font(r, size=FONT_SIZE_CAPTION, bold=False, color_hex=COLOR_MUTED)

    return pic


def add_bullets(slide, items: Sequence,
                top: Optional[float] = None,
                left: Optional[float] = None,
                width: Optional[float] = None,
                font_size: float = FONT_SIZE_SUBTITLE,
                title: Optional[str] = None):
    """
    加入項目列表。

    Args:
        items: list，每項可以是 str（第一層）或 tuple(text, indent_level)
        title: 列表上方粗體標題
    """
    spec = _get_slide_brand_spec(slide)
    _left = left if left is not None else spec.content_area.left
    _top = top if top is not None else spec.content_area.top
    _width = width if width is not None else spec.content_area.width

    # 估算高度：每項 0.35"，+ title 0.35"
    n = len(items) + (1 if title else 0)
    est_h = max(1.0, n * 0.38)

    box = slide.shapes.add_textbox(
        Inches(_left), Inches(_top),
        Inches(_width), Inches(est_h)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0

    first = True

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        _set_run_font(r, size=font_size + 2, bold=True, color_hex=COLOR_PRIMARY)
        p.space_after = Pt(6)
        first = False

    for item in items:
        if isinstance(item, tuple):
            text, indent = item
        else:
            text, indent = str(item), 1

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = max(0, int(indent) - 1)

        bullet = "• " if indent <= 1 else "– "
        r = p.add_run()
        r.text = bullet + text
        _set_run_font(r, size=font_size, bold=False, color_hex=COLOR_TEXT)
        p.space_after = Pt(4)

    return box


def add_text_box(slide, text: str,
                 left: float, top: float, width: float,
                 height: Optional[float] = None,
                 font_size: float = FONT_SIZE_BODY,
                 bold: bool = False,
                 color: str = COLOR_TEXT,
                 bg_color: Optional[str] = None,
                 border_left_color: Optional[str] = None):
    """
    通用文字框。border_left_color 非 None 時加左側色條（callout 效果）。
    """
    _height = height if height is not None else 0.8

    # 左側色條（橫向與文字框錯開，避免疊合）
    if border_left_color:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(0.08), Inches(_height)
        )
        _set_shape_fill(bar, border_left_color)
        _set_shape_no_line(bar)
        tb_left = left + 0.18
        tb_width = width - 0.18
    else:
        tb_left = left
        tb_width = width

    # 用 rectangle shape 同時承擔背景與文字（避免獨立 bg 矩形產生疊合）
    if bg_color:
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(tb_left), Inches(top),
            Inches(tb_width), Inches(_height)
        )
        _set_shape_fill(box, bg_color)
        _set_shape_no_line(box)
    else:
        box = slide.shapes.add_textbox(
            Inches(tb_left), Inches(top),
            Inches(tb_width), Inches(_height)
        )

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    _set_run_font(r, size=font_size, bold=bold, color_hex=color)

    return box


def add_two_column(slide,
                   left_content_fn: Callable,
                   right_content_fn: Callable,
                   ratio: str = "1:1",
                   top: Optional[float] = None,
                   gap: float = 0.25) -> None:
    """
    雙欄佈局。

    Args:
        left_content_fn / right_content_fn: callable(slide, left, top, width, content_height)
        ratio: "1:1" / "2:1" / "1:2"
        top: 欄起始 top（inches），預設 CONTENT_TOP
        gap: 欄間距（inches）
    """
    spec = _get_slide_brand_spec(slide)
    try:
        a_str, b_str = ratio.split(":")
        a, b = float(a_str), float(b_str)
    except Exception:
        a, b = 1.0, 1.0

    total = a + b
    _top = top if top is not None else spec.content_area.top
    content_h = (spec.content_area.top + spec.content_area.height) - _top

    left_w = (spec.content_area.width - gap) * (a / total)
    right_w = (spec.content_area.width - gap) * (b / total)

    left_x = spec.content_area.left
    right_x = spec.content_area.left + left_w + gap

    left_content_fn(slide, left_x, _top, left_w, content_h)
    right_content_fn(slide, right_x, _top, right_w, content_h)


def get_brand_safe_area(brand_spec: Optional[BrandSpec] = None) -> dict:
    """Module-level helper for callers that only need品牌母片安全區。"""
    spec = brand_spec or get_active_brand_spec()
    return {
        "content_left": spec.content_area.left,
        "content_top": spec.content_area.top,
        "content_width": spec.content_area.width,
        "content_height": spec.content_area.height,
        "content_bottom": spec.content_area.top + spec.content_area.height,
        "title_left": spec.title_area.left,
        "title_top": spec.title_area.top,
        "title_width": spec.title_area.width,
        "title_height": spec.title_area.height,
        "header_line_left": spec.header_line_area.left,
        "header_line_top": spec.header_line_area.top,
        "header_line_width": spec.header_line_area.width,
        "header_line_height": spec.header_line_area.height,
        "logo_left": spec.content_logo_area.left,
        "logo_top": spec.content_logo_area.top,
        "logo_width": spec.content_logo_area.width,
        "logo_height": spec.content_logo_area.height,
        "footer_left": spec.footer_area.left,
        "footer_top": spec.footer_area.top,
        "footer_width": spec.footer_area.width,
        "footer_height": spec.footer_area.height,
        "page_num_left": spec.page_num_area.left,
        "page_num_top": spec.page_num_area.top,
        "page_num_width": spec.page_num_area.width,
        "page_num_height": spec.page_num_area.height,
    }


# ============================================================================
# CLI — demo / info
# ============================================================================

PanjitPresentation = PanjitBrandMasterPresentation

def _cmd_demo(output_path: str) -> int:
    """產生示範 PPT：封面 + 章節 + 表格 + 雙欄 + callout + 結尾"""
    normalized_output = _normalize_cli_path(output_path)
    prs = PanjitBrandMasterPresentation()

    # Slide 1: 封面
    prs.add_title_slide(
        title="PANJIT 樣式 PPT 示範",
        subtitle="pptx_panjit.py 工具庫驗證報告",
        department="QA Engineering",
        owner="可靠性工程組",
        date="2026-04-21",
        version="V1.0",
    )

    # Slide 2: 章節分隔
    prs.add_section_divider("第一章　基礎元素展示")

    # Slide 3: 含表格
    slide = prs.add_master_content_slide(
        title="可靠性測試結果彙整",
        section_label="第一章 § 表格元素"
    )
    table_data = [
        ["測試項目", "規格", "樣品數", "Pass/Fail", "結果"],
        ["HTRL", "Ta=150°C, 1000h", "77", "0/77", "PASS"],
        ["HTGB", "VGS=-20V, 1000h", "77", "0/77", "PASS"],
        ["H3TRB", "85°C/85%RH, 1000h", "77", "0/77", "PASS"],
        ["TC", "-65°C ~ +150°C, 1000cyc", "77", "1/77", "FAIL"],
        ["PC", "-40°C ~ +125°C, 1000cyc", "77", "0/77", "PASS"],
    ]
    add_table(
        slide, table_data,
        col_widths=[2.8, 3.5, 1.5, 2.0, 1.8],
        caption="資料來源：AEC-Q101 Q1 2026 可靠性報告",
    )

    # Slide 4: 雙欄（左 bullets，右 callout）
    slide2 = prs.add_master_content_slide(
        title="風險評估與行動項目",
        section_label="第一章 § 雙欄佈局"
    )

    def _left_fn(s, left, top, width, h):
        add_bullets(
            s,
            items=[
                "TC 1000cyc 失效率 1.3%",
                ("疑似 die attach 裂紋", 2),
                ("需複測並做 SAM 分析", 2),
                "其他項目全數 PASS",
                "下一步：追加 500 樣品複測",
            ],
            left=left, top=top, width=width,
            title="關鍵發現",
            font_size=12,
        )

    def _right_fn(s, left, top, width, h):
        add_text_box(
            s,
            text=(
                "注意：TC 失效樣品已送至可靠性實驗室，"
                "預計 2026-05-10 前完成 SAM 掃描與 cross-section 分析。"
                "失效模式確認後將更新至 FMEA。"
            ),
            left=left, top=top, width=width,
            height=1.6,
            font_size=12,
            bg_color=COLOR_WARNING,
            border_left_color=COLOR_ACCENT,
        )
        add_text_box(
            s,
            text="負責人：林志明　預計完成：2026-05-10",
            left=left, top=top + 1.8, width=width,
            height=0.5,
            font_size=11,
            bold=True,
            color=COLOR_PRIMARY,
        )

    add_two_column(slide2, _left_fn, _right_fn, ratio="1:1")

    # Slide 5: 結尾頁
    slide3 = prs.add_master_content_slide(
        title="結論與下一步",
        section_label="結尾"
    )
    add_bullets(
        slide3,
        items=[
            "AEC-Q101 主要項目通過，TC 項需複測確認",
            "複測計畫：77pcs × 2 batch，2026-05 完成",
            "後續動作：",
            ("FMEA 更新（die attach 失效模式加權）", 2),
            ("供應商品質會議：2026-05-15", 2),
            ("正式報告 release：2026-05-30", 2),
        ],
        top=CONTENT_TOP + 0.1,
        font_size=14,
    )

    prs.save(normalized_output)
    print(f"OK: {normalized_output} ({Path(normalized_output).stat().st_size // 1024} KB)")

    # 跑 validator
    validator = Path(__file__).parent / "office_validator.py"
    if validator.exists():
        print("\n=== Running office_validator ===")
        result = subprocess.run(
            [sys.executable, str(validator), normalized_output],
            capture_output=True, text=True, encoding="utf-8", errors="replace",
        )
        print(result.stdout)
        if result.stderr:
            print("[stderr]", result.stderr, file=sys.stderr)
        return result.returncode
    return 0


def _cmd_demo_with_spec(spec_path: str, output_path: str) -> int:
    normalized_output = _normalize_cli_path(output_path)
    prs = PanjitBrandMasterPresentation.from_brand_spec_file(spec_path)
    slide = prs.add_title_slide(
        title="Brand Master Demo",
        subtitle=Path(spec_path).name,
        department="QA Engineering",
        owner="architect",
        date="2026-04-25",
        version="V1.0",
    )
    _ = slide
    body = prs.add_master_content_slide(
        title="Safe Area Preview",
        section_label="brand spec",
    )
    safe = prs.get_content_safe_area()
    add_text_box(
        body,
        text=json.dumps(safe, ensure_ascii=False, indent=2),
        left=safe["content"]["left"],
        top=safe["content"]["top"],
        width=min(6.8, safe["content"]["width"]),
        height=3.2,
        font_size=10,
        color=prs.brand_spec.color_text,
    )
    prs.save(normalized_output)
    print(f"OK: {normalized_output} ({Path(normalized_output).stat().st_size // 1024} KB)")
    return 0


def _cmd_info(template_path: str) -> int:
    """顯示 pptx 模板 layout 資訊"""
    from pptx import Presentation as _P
    normalized = _normalize_cli_path(template_path)
    prs = _P(normalized)
    print(f"File: {normalized}")
    print(f"Slide size: {Emu(prs.slide_width).inches:.2f}\" x {Emu(prs.slide_height).inches:.2f}\"")
    print(f"Slides: {len(prs.slides)}")
    print(f"Slide layouts ({len(prs.slide_layouts)}):")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")
    print(f"Slide masters ({len(prs.slide_masters)}):")
    for i, m in enumerate(prs.slide_masters):
        print(f"  [{i}] {m.name if hasattr(m, 'name') else '<master>'}")
    return 0


def _cmd_dump_spec(output_path: str) -> int:
    save_brand_spec(output_path, get_active_brand_spec())
    print(f"OK: wrote brand spec -> {output_path}")
    return 0


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            parts.append(text)
    return "\n".join(parts).strip()


def _rect_from_shape(shape) -> RectSpec:
    return RectSpec(
        left=Emu(shape.left).inches,
        top=Emu(shape.top).inches,
        width=Emu(shape.width).inches,
        height=Emu(shape.height).inches,
    )


def _emu_to_inches(value: Optional[str]) -> float:
    return int(value or 0) / 914400.0


def _resolve_zip_target(base_part: str, target: str) -> str:
    return posixpath.normpath(posixpath.join(posixpath.dirname(base_part), target))


def _extract_layout_artifacts_from_pptx(
    pptx_path: str,
    slide_number: int,
    logo_output: Optional[str] = None,
) -> dict:
    artifacts = {}
    slide_rels_part = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    with zipfile.ZipFile(pptx_path, "r") as zf:
        if slide_rels_part not in zf.namelist():
            return artifacts

        rels_root = ET.fromstring(zf.read(slide_rels_part))
        layout_target = None
        for rel in rels_root.findall("rel:Relationship", _PML_NS):
            rel_type = rel.get("Type", "")
            if rel_type.endswith("/slideLayout"):
                layout_target = rel.get("Target")
                break
        if not layout_target:
            return artifacts

        layout_part = _resolve_zip_target("ppt/slides/slide1.xml", layout_target)
        if layout_part not in zf.namelist():
            return artifacts

        layout_root = ET.fromstring(zf.read(layout_part))
        layout_rels_part = posixpath.join(
            posixpath.dirname(layout_part),
            "_rels",
            posixpath.basename(layout_part) + ".rels",
        )
        layout_rel_map = {}
        if layout_rels_part in zf.namelist():
            layout_rels_root = ET.fromstring(zf.read(layout_rels_part))
            for rel in layout_rels_root.findall("rel:Relationship", _PML_NS):
                layout_rel_map[rel.get("Id")] = rel.get("Target")

        for sp in layout_root.findall(".//p:sp", _PML_NS):
            off = sp.find(".//a:xfrm/a:off", _PML_NS)
            ext = sp.find(".//a:xfrm/a:ext", _PML_NS)
            rect = RectSpec(
                left=_emu_to_inches(off.get("x") if off is not None else None),
                top=_emu_to_inches(off.get("y") if off is not None else None),
                width=_emu_to_inches(ext.get("cx") if ext is not None else None),
                height=_emu_to_inches(ext.get("cy") if ext is not None else None),
            )
            text = " ".join(
                t.text.strip()
                for t in sp.findall(".//a:t", _PML_NS)
                if t.text and t.text.strip()
            ).strip()
            ph = sp.find(".//p:nvPr/p:ph", _PML_NS)
            ph_type = ph.get("type") if ph is not None else None
            fld = sp.find(".//a:fld", _PML_NS)
            fld_type = fld.get("type") if fld is not None else None

            if ph_type in {"title", "ctrTitle"} and rect.width > 0 and rect.height > 0:
                artifacts["title_area"] = rect
            elif ph_type == "body" and rect.width > 0 and rect.height > 0:
                artifacts["content_area"] = rect
            elif "Copyright" in text:
                artifacts["footer_area"] = rect
                artifacts["footer_text"] = text
            elif fld_type == "slidenum":
                artifacts["page_num_area"] = rect

        for cxn in layout_root.findall(".//p:cxnSp", _PML_NS):
            off = cxn.find(".//a:xfrm/a:off", _PML_NS)
            ext = cxn.find(".//a:xfrm/a:ext", _PML_NS)
            rect = RectSpec(
                left=_emu_to_inches(off.get("x") if off is not None else None),
                top=_emu_to_inches(off.get("y") if off is not None else None),
                width=_emu_to_inches(ext.get("cx") if ext is not None else None),
                height=_emu_to_inches(ext.get("cy") if ext is not None else None),
            )
            if rect.width >= 6.0 and rect.height <= 0.08 and rect.top <= 1.2:
                artifacts["header_line_area"] = rect
                break

        for pic in layout_root.findall(".//p:pic", _PML_NS):
            off = pic.find(".//a:xfrm/a:off", _PML_NS)
            ext = pic.find(".//a:xfrm/a:ext", _PML_NS)
            rect = RectSpec(
                left=_emu_to_inches(off.get("x") if off is not None else None),
                top=_emu_to_inches(off.get("y") if off is not None else None),
                width=_emu_to_inches(ext.get("cx") if ext is not None else None),
                height=_emu_to_inches(ext.get("cy") if ext is not None else None),
            )
            if rect.left >= 9.5 and rect.top <= 1.2:
                artifacts["content_logo_area"] = rect
                blip = pic.find(".//a:blip", _PML_NS)
                embed = blip.get(f"{{{_PML_NS['r']}}}embed") if blip is not None else None
                media_target = layout_rel_map.get(embed) if embed else None
                if media_target:
                    media_part = _resolve_zip_target(layout_part, media_target)
                    if media_part in zf.namelist():
                        ext_name = Path(media_part).suffix or ".png"
                        if logo_output:
                            logo_path = Path(_normalize_cli_path(logo_output))
                        else:
                            out_json = Path(_normalize_cli_path(pptx_path)).with_suffix("")
                            logo_path = out_json.with_name(out_json.name + "_logo" + ext_name)
                        logo_path.parent.mkdir(parents=True, exist_ok=True)
                        logo_path.write_bytes(zf.read(media_part))
                        artifacts["content_logo_path"] = str(logo_path)
                        artifacts["cover_logo_path"] = str(logo_path)
                break
    return artifacts


def _pick_content_slide(prs):
    if len(prs.slides) >= 2:
        return prs.slides[1]
    return prs.slides[0]


def _find_best_title_shape(slide):
    candidates = []
    for shape in slide.shapes:
        text = _shape_text(shape)
        if not text:
            continue
        rect = _rect_from_shape(shape)
        if rect.top <= 1.2 and rect.width >= 4.0:
            candidates.append((rect.top, -rect.width, shape))
    return candidates[0][2] if candidates else None


def _find_section_label_shape(slide, title_shape):
    if title_shape is None:
        return None
    title_rect = _rect_from_shape(title_shape)
    candidates = []
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        text = _shape_text(shape)
        if not text:
            continue
        rect = _rect_from_shape(shape)
        if title_rect.top < rect.top <= 1.5 and rect.height <= 0.5 and rect.width >= 2.0:
            candidates.append((rect.top, shape))
    return candidates[0][1] if candidates else None


def _find_footer_shape(slide, slide_height: float, slide_width: float):
    candidates = []
    for shape in slide.shapes:
        text = _shape_text(shape)
        if not text:
            continue
        rect = _rect_from_shape(shape)
        if rect.top >= slide_height - 0.5 and rect.left <= slide_width / 2:
            candidates.append((-len(text), rect.left, shape))
    return candidates[0][2] if candidates else None


def _find_page_num_shape(slide, slide_height: float, slide_width: float):
    candidates = []
    for shape in slide.shapes:
        text = _shape_text(shape)
        if not text:
            continue
        rect = _rect_from_shape(shape)
        if rect.top >= slide_height - 0.5 and rect.left >= slide_width * 0.75:
            candidates.append((rect.left, shape))
    return candidates[0][1] if candidates else None


def _find_header_line_shape(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            continue
        rect = _rect_from_shape(shape)
        if rect.top <= 1.2 and rect.height <= 0.08 and rect.width >= 6.0:
            candidates.append((-rect.width, rect.top, shape))
    return candidates[0][2] if candidates else None


def _find_logo_shape(slide, slide_width: float):
    candidates = []
    for shape in slide.shapes:
        rect = _rect_from_shape(shape)
        if rect.top <= 1.0 and rect.left >= slide_width * 0.7 and hasattr(shape, "image"):
            candidates.append((rect.top, rect.left, shape))
    return candidates[0][2] if candidates else None


def _extract_brand_spec_from_pptx(template_path: str, output_path: str, logo_output: Optional[str] = None) -> int:
    normalized = _normalize_cli_path(template_path)
    prs = Presentation(normalized)
    slide_width = Emu(prs.slide_width).inches
    slide_height = Emu(prs.slide_height).inches
    slide = _pick_content_slide(prs)

    spec = BrandSpec()
    spec.slide_width = slide_width
    spec.slide_height = slide_height

    title_shape = _find_best_title_shape(slide)
    if title_shape is not None:
        spec.title_area = _rect_from_shape(title_shape)

    section_shape = _find_section_label_shape(slide, title_shape)
    if section_shape is not None:
        spec.section_label_area = _rect_from_shape(section_shape)

    line_shape = _find_header_line_shape(slide)
    if line_shape is not None:
        spec.header_line_area = _rect_from_shape(line_shape)

    footer_shape = _find_footer_shape(slide, slide_height, slide_width)
    if footer_shape is not None:
        spec.footer_area = _rect_from_shape(footer_shape)
        footer_text = _shape_text(footer_shape)
        if footer_text:
            spec.footer_text = footer_text.replace("\n", " ")

    page_shape = _find_page_num_shape(slide, slide_height, slide_width)
    if page_shape is not None:
        spec.page_num_area = _rect_from_shape(page_shape)

    layout_artifacts = _extract_layout_artifacts_from_pptx(
        normalized,
        slide_number=2 if len(prs.slides) >= 2 else 1,
        logo_output=logo_output,
    )
    if "title_area" in layout_artifacts:
        spec.title_area = layout_artifacts["title_area"]
    if "content_area" in layout_artifacts:
        spec.content_area = layout_artifacts["content_area"]
    if "footer_area" in layout_artifacts:
        spec.footer_area = layout_artifacts["footer_area"]
    if "footer_text" in layout_artifacts:
        spec.footer_text = layout_artifacts["footer_text"]
    if "page_num_area" in layout_artifacts:
        spec.page_num_area = layout_artifacts["page_num_area"]
    if "header_line_area" in layout_artifacts:
        spec.header_line_area = layout_artifacts["header_line_area"]
    if "content_logo_area" in layout_artifacts:
        spec.content_logo_area = layout_artifacts["content_logo_area"]
    if "content_logo_path" in layout_artifacts:
        spec.content_logo_path = layout_artifacts["content_logo_path"]
    if "cover_logo_path" in layout_artifacts:
        spec.cover_logo_path = layout_artifacts["cover_logo_path"]

    logo_shape = _find_logo_shape(slide, slide_width)
    if logo_shape is not None and not spec.content_logo_path:
        spec.content_logo_area = _rect_from_shape(logo_shape)
        try:
            image = logo_shape.image
            ext = image.ext or "png"
            if logo_output:
                logo_path = Path(_normalize_cli_path(logo_output))
            else:
                out_json = Path(_normalize_cli_path(output_path))
                logo_path = out_json.with_name(out_json.stem + "_logo." + ext)
            logo_path.parent.mkdir(parents=True, exist_ok=True)
            logo_path.write_bytes(image.blob)
            spec.content_logo_path = str(logo_path)
            spec.cover_logo_path = str(logo_path)
        except Exception:
            pass

    if "content_area" not in layout_artifacts:
        content_top = spec.header_line_area.top + max(spec.header_line_area.height, 0.03) + 0.52
        content_bottom = spec.footer_area.top - 0.19
        spec.content_area = RectSpec(
            left=min(spec.title_area.left, CONTENT_LEFT),
            top=round(content_top, 2),
            width=round(spec.slide_width - min(spec.title_area.left, CONTENT_LEFT) - 0.41, 2),
            height=round(max(1.0, content_bottom - content_top), 2),
        )

    save_brand_spec(output_path, spec)
    print(f"OK: extracted brand spec -> {output_path}")
    if spec.content_logo_path:
        print(f"OK: extracted logo -> {spec.content_logo_path}")
    return 0


def main() -> int:
    parser = argparse.ArgumentParser(
        description="PANJIT 樣式 PPTX 工具（pptx_panjit）"
    )
    sub = parser.add_subparsers(dest="cmd", required=True)

    p_demo = sub.add_parser("demo", help="產生示範 PPT")
    p_demo.add_argument("output", help="輸出 .pptx 路徑")

    p_demo_spec = sub.add_parser("demo-with-spec", help="用 JSON brand spec 產生示範 PPT")
    p_demo_spec.add_argument("spec", help="brand spec JSON 路徑")
    p_demo_spec.add_argument("output", help="輸出 .pptx 路徑")

    p_info = sub.add_parser("info", help="顯示模板 layout 資訊")
    p_info.add_argument("template", help="輸入 .pptx 路徑")

    p_dump = sub.add_parser("dump-spec", help="匯出預設 brand spec JSON")
    p_dump.add_argument("output", help="輸出 JSON 路徑")

    p_extract = sub.add_parser("extract-spec", help="從既有 PPTX 範例萃取 brand spec JSON")
    p_extract.add_argument("template", help="輸入 .pptx 路徑")
    p_extract.add_argument("output", help="輸出 JSON 路徑")
    p_extract.add_argument("--logo-output", dest="logo_output", help="輸出 logo 圖片路徑")

    args = parser.parse_args()

    if args.cmd == "demo":
        return _cmd_demo(args.output)
    if args.cmd == "demo-with-spec":
        return _cmd_demo_with_spec(args.spec, args.output)
    if args.cmd == "info":
        return _cmd_info(args.template)
    if args.cmd == "dump-spec":
        return _cmd_dump_spec(args.output)
    if args.cmd == "extract-spec":
        return _extract_brand_spec_from_pptx(args.template, args.output, args.logo_output)
    return 1


if __name__ == "__main__":
    sys.exit(main())
