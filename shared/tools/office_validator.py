"""
Office Document Validator — 驗證 PPTX / DOCX / XLSX 產出品質

用法：
    python shared/tools/office_validator.py <file_path> [--json] [--fix]

驗證項目：
  PPTX: 元素重疊、超出邊界、文字太小、空白投影片、殘留 placeholder
  DOCX: 圖片/表格超出頁寬、文字太小、缺少頁首頁尾、標題層級
  XLSX: 欄位過寬/過窄、凍結窗格、header 格式

回傳 exit code：0=PASS, 1=有 WARNING, 2=有 ERROR
"""

import argparse
import json
import re
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Optional

# Avoid local `shared/tools/pptx/` shadowing the third-party `python-pptx` package.
_THIS_DIR = Path(__file__).resolve().parent
sys.path = [p for p in sys.path if Path(p).resolve() != _THIS_DIR]


def _normalize_cli_path(path: str) -> str:
    if not path:
        return path
    m = re.match(r"^/([a-zA-Z])/(.+)$", path)
    if m:
        drive = m.group(1).upper()
        tail = m.group(2).replace("/", "\\")
        return f"{drive}:\\{tail}"
    return path


# ---------------------------------------------------------------------------
# 資料結構
# ---------------------------------------------------------------------------

@dataclass
class Issue:
    severity: str          # ERROR / WARNING / INFO
    rule: str              # 規則代碼 e.g. "PPTX-001"
    message: str           # 人類可讀描述
    location: str = ""     # e.g. "Slide 3, Shape 'Title 1'"
    suggestion: str = ""   # 修正建議


@dataclass
class ValidationResult:
    file_path: str
    file_type: str
    issues: list[Issue] = field(default_factory=list)
    summary: dict = field(default_factory=dict)

    @property
    def verdict(self) -> str:
        severities = {i.severity for i in self.issues}
        if "ERROR" in severities:
            return "FAIL"
        if "WARNING" in severities:
            return "WARNING"
        return "PASS"

    def add(self, severity: str, rule: str, message: str,
            location: str = "", suggestion: str = ""):
        self.issues.append(Issue(severity, rule, message, location, suggestion))

    def to_text(self) -> str:
        lines = [f"=== Office Validator: {self.file_type} ===",
                 f"File: {self.file_path}",
                 f"Verdict: {self.verdict}",
                 ""]

        if self.summary:
            lines.append("--- Summary ---")
            for k, v in self.summary.items():
                lines.append(f"  {k}: {v}")
            lines.append("")

        if not self.issues:
            lines.append("No issues found.")
            return "\n".join(lines)

        # Group by severity
        for sev in ("ERROR", "WARNING", "INFO"):
            group = [i for i in self.issues if i.severity == sev]
            if not group:
                continue
            lines.append(f"--- {sev} ({len(group)}) ---")
            for i in group:
                loc = f" @ {i.location}" if i.location else ""
                lines.append(f"  [{i.rule}]{loc}: {i.message}")
                if i.suggestion:
                    lines.append(f"    → {i.suggestion}")
            lines.append("")

        return "\n".join(lines)

    def to_dict(self) -> dict:
        return {
            "file_path": self.file_path,
            "file_type": self.file_type,
            "verdict": self.verdict,
            "summary": self.summary,
            "issues": [asdict(i) for i in self.issues],
        }


# ---------------------------------------------------------------------------
# PPTX Validator
# ---------------------------------------------------------------------------

# 標準寬螢幕投影片尺寸（EMU）
SLIDE_WIDTH_STANDARD = 12192000   # 13.333"
SLIDE_HEIGHT_STANDARD = 6858000   # 7.5"

# 閾值
MIN_FONT_SIZE_PT = 10
OVERLAP_THRESHOLD_EMU = 18000  # ~0.02" 容差，避免邊緣剛好觸碰的假陽性


def _emu_to_inches(emu: int) -> float:
    return emu / 914400


def _rects_overlap(s1, s2, threshold: int = OVERLAP_THRESHOLD_EMU) -> bool:
    """檢查兩個 shape 的 bounding box 是否重疊（含容差）。"""
    l1, t1, r1, b1 = s1.left, s1.top, s1.left + s1.width, s1.top + s1.height
    l2, t2, r2, b2 = s2.left, s2.top, s2.left + s2.width, s2.top + s2.height
    # 縮小各邊 threshold 來避免邊緣觸碰的假陽性
    return (l1 + threshold < r2 and r1 - threshold > l2 and
            t1 + threshold < b2 and b1 - threshold > t2)


def _overlap_area_ratio(s1, s2) -> float:
    """計算重疊面積佔較小 shape 面積的比例。"""
    l1, t1, r1, b1 = s1.left, s1.top, s1.left + s1.width, s1.top + s1.height
    l2, t2, r2, b2 = s2.left, s2.top, s2.left + s2.width, s2.top + s2.height
    ol = max(l1, l2)
    ot = max(t1, t2)
    or_ = min(r1, r2)
    ob = min(b1, b2)
    if ol >= or_ or ot >= ob:
        return 0.0
    overlap_area = (or_ - ol) * (ob - ot)
    area1 = max(s1.width * s1.height, 1)
    area2 = max(s2.width * s2.height, 1)
    return overlap_area / min(area1, area2)


def validate_pptx(file_path: str) -> ValidationResult:
    from pptx import Presentation
    from pptx.util import Pt

    result = ValidationResult(file_path=file_path, file_type="PPTX")
    prs = Presentation(file_path)

    slide_w = prs.slide_width or SLIDE_WIDTH_STANDARD
    slide_h = prs.slide_height or SLIDE_HEIGHT_STANDARD

    result.summary = {
        "slide_count": len(prs.slides),
        "slide_size": f"{_emu_to_inches(slide_w):.2f}\" x {_emu_to_inches(slide_h):.2f}\"",
    }

    for slide_idx, slide in enumerate(prs.slides, 1):
        loc_slide = f"Slide {slide_idx}"
        shapes = list(slide.shapes)

        # --- PPTX-001: 空白投影片 ---
        if not shapes:
            result.add("WARNING", "PPTX-001",
                       "Empty slide (no shapes)",
                       loc_slide,
                       "Add content or remove the slide")
            continue

        # --- PPTX-002: 殘留預設 placeholder ---
        for shape in shapes:
            try:
                ph_fmt = shape.placeholder_format
            except ValueError:
                ph_fmt = None
            if ph_fmt is not None:
                ph_type = ph_fmt.type
                # Type 13=TITLE, 14=SUBTITLE in blank layouts
                if ph_type is not None and int(ph_type) in (13, 14):
                    text = shape.text_frame.text.strip() if shape.has_text_frame else ""
                    if not text:
                        result.add("WARNING", "PPTX-002",
                                   f"Unused placeholder left: {shape.name} (type={ph_type})",
                                   f"{loc_slide}, Shape '{shape.name}'",
                                   "Delete with delete_shape or add content")

        # --- PPTX-003: 元素超出投影片邊界 ---
        for shape in shapes:
            right = shape.left + shape.width
            bottom = shape.top + shape.height

            violations = []
            if shape.left < 0:
                violations.append(f"left edge at {_emu_to_inches(shape.left):.2f}\" (< 0)")
            if shape.top < 0:
                violations.append(f"top edge at {_emu_to_inches(shape.top):.2f}\" (< 0)")
            if right > slide_w:
                overflow = _emu_to_inches(right - slide_w)
                violations.append(f"right edge exceeds by {overflow:.2f}\"")
            if bottom > slide_h:
                overflow = _emu_to_inches(bottom - slide_h)
                violations.append(f"bottom edge exceeds by {overflow:.2f}\"")

            if violations:
                result.add("ERROR", "PPTX-003",
                           f"Shape exceeds slide boundary: {'; '.join(violations)}",
                           f"{loc_slide}, Shape '{shape.name}' "
                           f"(pos={_emu_to_inches(shape.left):.2f}\",{_emu_to_inches(shape.top):.2f}\" "
                           f"size={_emu_to_inches(shape.width):.2f}\"x{_emu_to_inches(shape.height):.2f}\")",
                           f"Adjust position/size to fit within {_emu_to_inches(slide_w):.1f}\"x{_emu_to_inches(slide_h):.1f}\"")

        # --- PPTX-004: 元素重疊 ---
        for i, s1 in enumerate(shapes):
            for s2 in shapes[i+1:]:
                if _rects_overlap(s1, s2):
                    ratio = _overlap_area_ratio(s1, s2)
                    if ratio > 0.5:
                        sev = "ERROR"
                        msg = f"Major overlap ({ratio:.0%})"
                    elif ratio > 0.1:
                        sev = "WARNING"
                        msg = f"Partial overlap ({ratio:.0%})"
                    else:
                        continue  # 微量重疊（< 10%）忽略

                    result.add(sev, "PPTX-004",
                               f"{msg} between '{s1.name}' and '{s2.name}'",
                               loc_slide,
                               "Reposition shapes to avoid overlap")

        # --- PPTX-005: 文字太小（每 shape 只報一次，摘要最小字級和數量）---
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            small_runs = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None and run.font.size.pt < MIN_FONT_SIZE_PT:
                        small_runs.append(run.font.size.pt)
            if small_runs:
                min_size = min(small_runs)
                result.add("WARNING", "PPTX-005",
                           f"{len(small_runs)} text runs below {MIN_FONT_SIZE_PT}pt "
                           f"(smallest: {min_size}pt)",
                           f"{loc_slide}, Shape '{shape.name}'",
                           f"Increase font size to at least {MIN_FONT_SIZE_PT}pt")

        # --- PPTX-006: 文字溢出文字框（text frame 設定 auto_size=None 時可能截斷）---
        for shape in shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            total_text = sum(len(run.text) for para in tf.paragraphs for run in para.runs)
            if total_text == 0:
                continue
            # 粗估：shape 面積 vs 文字量（每字約需 0.04 sq inches @14pt）
            area_sq_in = _emu_to_inches(shape.width) * _emu_to_inches(shape.height)
            avg_font = 14  # 預設估計
            sizes = [run.font.size.pt for para in tf.paragraphs for run in para.runs
                     if run.font.size is not None]
            if sizes:
                avg_font = sum(sizes) / len(sizes)
            # 每字大約需要 (font_pt/72)^2 * 0.6 平方英吋
            char_area = (avg_font / 72) ** 2 * 0.6
            needed_area = total_text * char_area
            if needed_area > area_sq_in * 1.2:  # 容差 20%
                result.add("WARNING", "PPTX-006",
                           f"Text may overflow: ~{total_text} chars in {area_sq_in:.1f} sq\" "
                           f"(estimated need {needed_area:.1f} sq\" at {avg_font:.0f}pt)",
                           f"{loc_slide}, Shape '{shape.name}'",
                           "Enlarge the text box or reduce text/font size")

    return result


# ---------------------------------------------------------------------------
# DOCX Validator
# ---------------------------------------------------------------------------

def validate_docx(file_path: str) -> ValidationResult:
    from docx import Document
    from docx.shared import Inches, Pt, Emu

    result = ValidationResult(file_path=file_path, file_type="DOCX")
    doc = Document(file_path)

    # 頁面尺寸
    section = doc.sections[0] if doc.sections else None
    if section:
        page_w = section.page_width
        left_m = section.left_margin
        right_m = section.right_margin
        content_w = page_w - left_m - right_m
        content_w_in = _emu_to_inches(content_w)
    else:
        content_w = Inches(6.5)
        content_w_in = 6.5

    result.summary = {
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
        "section_count": len(doc.sections),
        "content_width": f"{content_w_in:.2f}\"",
    }

    # --- DOCX-001: 文字太小 ---
    for para_idx, para in enumerate(doc.paragraphs, 1):
        for run in para.runs:
            if run.font.size is not None and run.font.size.pt < MIN_FONT_SIZE_PT:
                text_preview = run.text[:30]
                result.add("WARNING", "DOCX-001",
                           f"Font size {run.font.size.pt}pt < {MIN_FONT_SIZE_PT}pt: \"{text_preview}\"",
                           f"Paragraph {para_idx}",
                           f"Increase to at least {MIN_FONT_SIZE_PT}pt")

    # --- DOCX-002: 表格寬度檢查 ---
    for tbl_idx, table in enumerate(doc.tables, 1):
        if not table.rows:
            continue
        # 嘗試計算表格總寬
        first_row = table.rows[0]
        total_w = 0
        measurable = True
        for cell in first_row.cells:
            if cell.width:
                total_w += cell.width
            else:
                measurable = False
                break
        if measurable and total_w > content_w * 1.05:  # 5% 容差
            overflow_in = _emu_to_inches(total_w - content_w)
            result.add("ERROR", "DOCX-002",
                       f"Table width ({_emu_to_inches(total_w):.2f}\") exceeds "
                       f"content area ({content_w_in:.2f}\") by {overflow_in:.2f}\"",
                       f"Table {tbl_idx}",
                       "Reduce column widths or switch to landscape orientation")

    # --- DOCX-003: 圖片寬度檢查 ---
    from docx.oxml.ns import qn
    for para_idx, para in enumerate(doc.paragraphs, 1):
        for run in para.runs:
            drawings = run._element.findall(qn('w:drawing'))
            for drawing in drawings:
                # 取得 extent（圖片顯示尺寸）
                extents = drawing.findall('.//' + qn('wp:extent'))
                for ext in extents:
                    cx = int(ext.get('cx', 0))
                    if cx > content_w * 1.05:
                        overflow_in = _emu_to_inches(cx - content_w)
                        result.add("ERROR", "DOCX-003",
                                   f"Image width ({_emu_to_inches(cx):.2f}\") exceeds "
                                   f"content area ({content_w_in:.2f}\") by {overflow_in:.2f}\"",
                                   f"Paragraph {para_idx}",
                                   "Resize image to fit within page margins")

    # --- DOCX-004: 缺少頁首/頁尾/頁碼 ---
    if section:
        has_header = False
        has_footer = False
        for sec in doc.sections:
            if sec.header and sec.header.paragraphs:
                header_text = "".join(p.text for p in sec.header.paragraphs).strip()
                if header_text:
                    has_header = True
            if sec.footer and sec.footer.paragraphs:
                footer_text = "".join(p.text for p in sec.footer.paragraphs).strip()
                if footer_text:
                    has_footer = True
                # 檢查頁碼 field code
                for p in sec.footer.paragraphs:
                    if p._element.findall('.//' + qn('w:fldChar')):
                        has_footer = True

        if not has_header:
            result.add("INFO", "DOCX-004",
                       "No header found in document",
                       "Document sections",
                       "Add header with add_header_footer(header_text=...)")
        if not has_footer:
            result.add("INFO", "DOCX-004",
                       "No footer/page number found",
                       "Document sections",
                       "Add footer with add_header_footer(page_number=true)")

    # --- DOCX-005: 標題層級跳躍 ---
    heading_levels = []
    for para in doc.paragraphs:
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            try:
                level = int(para.style.name.replace("Heading ", "").strip())
                heading_levels.append((level, para.text[:40]))
            except ValueError:
                pass

    for i in range(1, len(heading_levels)):
        prev_level = heading_levels[i-1][0]
        curr_level = heading_levels[i][0]
        if curr_level > prev_level + 1:
            result.add("WARNING", "DOCX-005",
                       f"Heading level jumps from H{prev_level} to H{curr_level} "
                       f"(skipping H{prev_level+1})",
                       f"Heading: \"{heading_levels[i][1]}\"",
                       f"Use Heading {prev_level+1} or restructure hierarchy")

    # --- DOCX-006: 表格 header 樣式檢查 ---
    for tbl_idx, table in enumerate(doc.tables, 1):
        if not table.rows:
            continue
        header_row = table.rows[0]
        has_styled_header = False
        for cell in header_row.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    if run.bold or (run.font.color and run.font.color.rgb):
                        has_styled_header = True
                        break
            # 也檢查 cell shading
            tc = cell._tc
            shading = tc.findall('.//' + qn('w:shd'))
            if shading:
                has_styled_header = True
            if has_styled_header:
                break

        if not has_styled_header and len(table.rows) > 1:
            result.add("INFO", "DOCX-006",
                       "Table header row has no distinct styling (no bold/color/shading)",
                       f"Table {tbl_idx}",
                       "Apply header formatting: bold text + background color")

    return result


# ---------------------------------------------------------------------------
# XLSX Validator
# ---------------------------------------------------------------------------

def validate_xlsx(file_path: str) -> ValidationResult:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    result = ValidationResult(file_path=file_path, file_type="XLSX")
    wb = load_workbook(file_path, read_only=False, data_only=True)

    result.summary = {
        "sheet_count": len(wb.sheetnames),
        "sheets": wb.sheetnames,
    }

    for ws_name in wb.sheetnames:
        ws = wb[ws_name]
        loc_sheet = f"Sheet '{ws_name}'"

        # 跳過空 sheet
        if ws.max_row is None or ws.max_row < 1:
            result.add("WARNING", "XLSX-001",
                       "Empty worksheet",
                       loc_sheet,
                       "Add data or remove the sheet")
            continue

        # --- XLSX-002: 凍結窗格 ---
        if ws.max_row > 5 and ws.freeze_panes is None:
            result.add("INFO", "XLSX-002",
                       f"No freeze panes set ({ws.max_row} rows)",
                       loc_sheet,
                       "Freeze top row: freeze_panes(row=2, column=1)")

        # --- XLSX-003: 欄位寬度異常 ---
        for col_idx in range(1, min((ws.max_column or 1) + 1, 100)):
            col_letter = get_column_letter(col_idx)
            dim = ws.column_dimensions.get(col_letter)
            if dim and dim.width is not None:
                if dim.width > 80:
                    result.add("WARNING", "XLSX-003",
                               f"Column {col_letter} width={dim.width:.1f} (very wide)",
                               loc_sheet,
                               "Use auto_fit_columns or set reasonable width")
                elif dim.width < 2 and not dim.hidden:
                    result.add("WARNING", "XLSX-003",
                               f"Column {col_letter} width={dim.width:.1f} (too narrow to read)",
                               loc_sheet,
                               "Increase column width or use auto_fit_columns")

        # --- XLSX-004: header 列字型/格式檢查 ---
        if ws.max_row and ws.max_row > 1:
            header_styled = False
            for col_idx in range(1, min((ws.max_column or 1) + 1, 50)):
                cell = ws.cell(row=1, column=col_idx)
                if cell.value is not None:
                    if (cell.font and cell.font.bold) or \
                       (cell.fill and cell.fill.fgColor and
                        cell.fill.fgColor.rgb and cell.fill.fgColor.rgb != "00000000"):
                        header_styled = True
                        break
            if not header_styled:
                result.add("INFO", "XLSX-004",
                           "Header row (row 1) has no distinct styling",
                           loc_sheet,
                           "Apply header style: apply_style_preset('header', ...)")

        # --- XLSX-005: 文字太小 ---
        sample_rows = min(ws.max_row or 1, 20)
        sample_cols = min((ws.max_column or 1), 30)
        for row_idx in range(1, sample_rows + 1):
            for col_idx in range(1, sample_cols + 1):
                cell = ws.cell(row=row_idx, column=col_idx)
                if cell.font and cell.font.size and cell.font.size < 8:
                    result.add("WARNING", "XLSX-005",
                               f"Font size {cell.font.size}pt at cell "
                               f"{get_column_letter(col_idx)}{row_idx}",
                               loc_sheet,
                               "Increase font size to at least 8pt")
                    break  # 只報一次 per row
            else:
                continue
            break  # 只報一次 per sheet

        # --- XLSX-006: 自動篩選檢查 ---
        if ws.max_row and ws.max_row > 10 and ws.auto_filter.ref is None:
            result.add("INFO", "XLSX-006",
                       f"No auto filter set ({ws.max_row} rows)",
                       loc_sheet,
                       "Add auto filter: auto_filter(...)")

    wb.close()
    return result


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------

VALIDATORS = {
    ".pptx": validate_pptx,
    ".docx": validate_docx,
    ".xlsx": validate_xlsx,
}


def validate(file_path: str) -> ValidationResult:
    """驗證 Office 文件，根據副檔名自動選擇驗證器。"""
    normalized = _normalize_cli_path(file_path)
    p = Path(normalized)
    if not p.exists():
        r = ValidationResult(file_path=normalized, file_type="UNKNOWN")
        r.add("ERROR", "FILE-001", f"File not found: {normalized}")
        return r

    ext = p.suffix.lower()
    validator = VALIDATORS.get(ext)
    if not validator:
        r = ValidationResult(file_path=normalized, file_type=ext)
        r.add("ERROR", "FILE-002", f"Unsupported file type: {ext}")
        return r

    return validator(normalized)


# ---------------------------------------------------------------------------
# Visual QA: PDF → PNG 轉換（供 Claude Vision 審查）
# ---------------------------------------------------------------------------

def render_to_images(file_path: str, output_dir: Optional[str] = None,
                     scale: float = 2.0) -> list[str]:
    """將 Office 文件（需先轉為 PDF）或 PDF 直接渲染為 PNG 圖片。

    Args:
        file_path: .pdf 檔案路徑
        output_dir: 輸出目錄（預設為 PDF 同目錄下的 _visual_qa/）
        scale: 渲染倍率（2.0 = ~144 DPI，適合視覺審查）

    Returns:
        產出的 PNG 檔案路徑列表
    """
    import pypdfium2 as pdfium

    normalized = _normalize_cli_path(file_path)
    p = Path(normalized)
    if not p.exists():
        raise FileNotFoundError(f"PDF not found: {normalized}")

    if output_dir is None:
        output_dir = str(p.parent / "_visual_qa")
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    pdf = pdfium.PdfDocument(str(p))
    png_paths = []
    for page_idx in range(len(pdf)):
        page = pdf[page_idx]
        bitmap = page.render(scale=scale)
        pil_image = bitmap.to_pil()
        png_name = f"{p.stem}_page{page_idx + 1:02d}.png"
        png_path = str(Path(output_dir) / png_name)
        pil_image.save(png_path)
        png_paths.append(png_path)
    pdf.close()

    return png_paths


VISUAL_QA_PROMPT = """Visually inspect this document page. Assume there are issues — find them.

Look for:
- Overlapping elements (text through shapes, lines through words, stacked elements)
- Text overflow or cut off at edges/box boundaries
- Elements too close (< 0.3" gaps) or sections nearly touching
- Uneven gaps (large empty area in one place, cramped in another)
- Insufficient margin from page/slide edges (< 0.5")
- Columns or similar elements not aligned consistently
- Low-contrast text (light text on light background, dark on dark)
- Text too small to read comfortably
- Tables or images extending beyond page boundaries
- Leftover placeholder content or template artifacts
- Font inconsistency (mixed fonts within same hierarchy level)

For each page, list ALL issues found, even minor ones.
If no issues found, state "No issues" — but look again more critically first.
"""


def visual_qa(file_path: str, output_dir: Optional[str] = None) -> dict:
    """執行 Visual QA：Office/PDF → PNG，回傳圖片路徑和檢查提示。

    回傳格式：
    {
        "pdf_path": "...",
        "png_paths": ["page01.png", ...],
        "page_count": N,
        "prompt": "... (visual inspection prompt)",
        "instructions": "Use Claude Read tool to view each PNG, then report issues."
    }
    """
    normalized = _normalize_cli_path(file_path)
    p = Path(normalized)
    ext = p.suffix.lower()

    # 若非 PDF，需先轉換（提示使用者用 MCP 匯出）
    if ext != ".pdf":
        return {
            "error": f"File is {ext}, not PDF. Export to PDF first.",
            "instructions": {
                ".docx": "Use mcp__docx__export_to_pdf(output_path='...pdf')",
                ".pptx": "Use mcp__pptx__export_to_pdf(output_path='...pdf')",
                ".xlsx": "Print to PDF via Excel, or use openpyxl for structure validation only",
            }.get(ext, "Convert to PDF first."),
        }

    png_paths = render_to_images(normalized, output_dir)

    return {
        "pdf_path": str(p),
        "png_paths": png_paths,
        "page_count": len(png_paths),
        "prompt": VISUAL_QA_PROMPT,
        "instructions": (
            f"Rendered {len(png_paths)} pages to PNG.\n"
            "Use Claude Read tool to view each PNG file, then apply the prompt above.\n"
            "Report ALL issues found per page."
        ),
    }


def main():
    parser = argparse.ArgumentParser(description="Office Document Validator")
    parser.add_argument("file_path", help="Path to the Office file to validate")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    parser.add_argument("--visual", action="store_true",
                        help="Visual QA mode: render PDF to PNG for vision inspection")
    parser.add_argument("--output-dir", default=None,
                        help="Output directory for visual QA PNGs")
    args = parser.parse_args()

    if args.visual:
        result = visual_qa(args.file_path, args.output_dir)
        print(json.dumps(result, ensure_ascii=False, indent=2))
        sys.exit(0 if "error" not in result else 1)

    result = validate(args.file_path)

    if args.json:
        print(json.dumps(result.to_dict(), ensure_ascii=False, indent=2))
    else:
        print(result.to_text())

    # Exit code
    if result.verdict == "FAIL":
        sys.exit(2)
    elif result.verdict == "WARNING":
        sys.exit(1)
    else:
        sys.exit(0)


if __name__ == "__main__":
    main()
